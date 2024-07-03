import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches

def extract_key_frames(video_path, threshold=30):
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        print("Error: Could not open video.")
        return []

    frames = []
    ret, prev_frame = cap.read()
    if not ret:
        print("Error: Could not read the first frame.")
        return []

    frames.append(prev_frame)
    while True:
        ret, frame = cap.read()
        if not ret:
            break

        diff = cv2.absdiff(cv2.cvtColor(prev_frame, cv2.COLOR_BGR2GRAY), cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY))
        non_zero_count = np.count_nonzero(diff)
        if non_zero_count > threshold:
            frames.append(frame)
            prev_frame = frame

    cap.release()
    return frames

def create_ppt_from_frames(frames, output_ppt_path):
    prs = Presentation()
    for i, frame in enumerate(frames):
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        image_path = f"frame_{i}.jpg"
        cv2.imwrite(image_path, frame)
        slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))

    prs.save(output_ppt_path)

if __name__ == "__main__":
    video_path = "your_video.mp4"
    output_ppt_path = "output_presentation.pptx"
    frames = extract_key_frames(video_path)
    create_ppt_from_frames(frames, output_ppt_path)
    print(f"PowerPoint presentation saved as {output_ppt_path}")